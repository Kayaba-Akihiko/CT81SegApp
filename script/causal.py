

from pathlib import Path
import numpy as np
import io

import imageio.v3 as iio
from pptx import Presentation
from pptx.util import Inches

from lightning.fabric.strategies.launchers.multiprocessing import _MultiProcessingLauncher
from lightning.fabric.strategies.launchers.subprocess_script import _SubprocessScriptLauncher


def main():

    this_file = Path(__file__)
    this_dir = this_file.parent
    output_dir = this_dir / f'out_{this_file.stem}'
    output_dir.mkdir(exist_ok=True, parents=True)

    # 1. Create a sample NumPy array representing an image (e.g., a simple gradient)
    #    For a real image, you would load it into a NumPy array.
    width, height = 640, 480
    data = np.zeros((height, width, 3), dtype=np.uint8)
    for y in range(height):
        for x in range(width):
            data[y, x] = [int(255 * (x / width)), int(255 * (y / height)), 128]

    # 3. Save the Pillow Image to an in-memory bytes stream
    image_stream = io.BytesIO()
    iio.imwrite(image_stream, data, extension='.png')
    image_stream.seek(0)  # Reset stream position to the beginning

    # 4. Create a presentation and add a slide
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Index 6 is often a blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # 5. Add the image from the stream to the slide
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(3.75)  # Maintain aspect ratio if desired
    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)

    # 6. Save the presentation
    prs.save(str(output_dir / 'temp.pptx'))


if __name__ == '__main__':
    main()