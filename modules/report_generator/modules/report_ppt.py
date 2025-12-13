#  Copyright (c) 2025. by Yi GU <gu.yi.gu4@naist.ac.jp>,
#  Biomedical Imaging Intelligence Laboratory,
#  Nara Institute of Science and Technology.
#  All rights reserved.
#  This file can not be copied and/or distributed
#  without the express permission of Yi GU.


from pathlib import Path
from typing import Union, Dict, TypeAlias, Literal, Optional, Tuple, Self, IO
import io
from dataclasses import dataclass
import logging
import tempfile
import subprocess
import shutil
import imageio.v3 as iio

import numpy as np
import numpy.typing as npt
from pptx.presentation import (
    Presentation as PPTXPresentation,
)
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.shapes.base import BaseShape as PPTXBaseShape
from pptx.util import Length as PPTXLength
from pptx.text.text import _Run as PPTXRun
from pptx.slide import (
    Slide as PPTXSlide,
    SlideShapes as PPTXSlideShapes
)
from pptx import Presentation as create_presentation

from xmodules.xutils import os_utils
from xmodules.typing import TypePathLike

TypeFitMode: TypeAlias = Literal[
    'contain', 'cover', 'match_height', 'match_width']

_logger = logging.getLogger(__name__)

@dataclass
class FillImageData:
    name: Optional[str] = None
    image: Optional[npt.NDArray[np.uint8]] = None
    fill_mode: TypeFitMode = 'contain'
    send_to_back: bool = False


class ReportPPT:
    def __init__(
            self,
            template: Union[PPTXPresentation, Path],
            fit_mode: Union[TypeFitMode, Dict[str, TypeFitMode]] = 'contain',
    ):
        if isinstance(fit_mode, str):
            assert fit_mode in {'contain', 'cover', 'match_height', 'match_width'}
        elif isinstance(fit_mode, dict):
            for k, v in fit_mode.items():
                assert v in {'contain', 'cover', 'match_height', 'match_width'}
        else:
            raise ValueError(f'Invalid fit mode: {fit_mode}')
        self.fit_mode = fit_mode
        if isinstance(template, PPTXPresentation):
            self.presentation = self.clone_presentation(template)
        else:
            self.presentation = create_presentation(str(template))

    def copy(self) -> Self:
        return ReportPPT(self.presentation, fit_mode=self.fit_mode)

    @staticmethod
    def _extract_img_key(shape: PPTXBaseShape) -> Optional[str]:
        """Return IMG key from shape.name like '...IMG:foo', else None."""
        name = (getattr(shape, "name", "") or "")
        marker = "IMG:"
        idx = name.find(marker)
        if idx < 0:
            return None
        key = name[idx + len(marker):].strip()
        return key or None

    def collect_shape_sizes(self, dpi: int = 96) -> Dict[str, Tuple[int, int]]:
        """
        Collect placeholder (IMG:key) sizes in pixels.

        If multiple shapes share the same key, the last one wins (and a warning is logged).
        """
        sizes: Dict[str, Tuple[int, int]] = {}

        for slide in self.presentation.slides:
            for shape in self._iter_shapes(slide.shapes):
                key = self._extract_img_key(shape)
                if not key:
                    continue

                w_px = self._emu_to_px(shape.width, dpi)
                h_px = self._emu_to_px(shape.height, dpi)

                if key in sizes:
                    _logger.warning(f"Duplicate IMG key '{key}' found; overriding previous size.")
                sizes[key] = (w_px, h_px)

        return sizes

    def fill_images(self, images: Dict[str, FillImageData]) -> Dict[str, PPTXPicture]:
        """
        Fill IMG:key placeholders with images.

        - Validates missing images early
        - Returns dict key->picture
        """
        res: Dict[str, "PPTXPicture"] = {}
        finished_keys: set[str] = set()

        # Pre-validate input (fail fast)
        for k, d in images.items():
            if d.image is None:
                raise ValueError(f"Image not found: {k}")

        for slide in self.presentation.slides:
            for shape in self._iter_shapes(slide.shapes):
                key = self._extract_img_key(shape)
                if not key:
                    continue
                if key not in images:
                    continue

                d = images[key]
                pic = self._fit_image_into_slide(
                    slide,
                    d.image,
                    (shape.left, shape.top, shape.width, shape.height),
                    d.fill_mode,
                )

                if d.send_to_back:
                    # Best-effort: move element earlier within its parent
                    el = pic._element
                    parent = el.getparent()
                    if parent is not None:
                        try:
                            parent.remove(el)
                            parent.insert(0, el)  # "back" (usually)
                        except Exception:
                            _logger.exception("Failed to send picture to back for key=%s", key)

                res[key] = pic
                finished_keys.add(key)

        unknown = set(images.keys()) - finished_keys
        if unknown:
            _logger.warning(f"Unknown filling keys: {sorted(unknown)}")

        return res

    @staticmethod
    def _replace_tokens_in_runs(
            runs, token_map: Dict[str, str], res: Dict[str, PPTXRun], finished: set[str]) -> None:
        """
        Replace tokens in a list of runs in-place.
        token_map: token -> replacement string
        res: mapping from key -> run (first run where key was replaced, best-effort)
        """
        # To allow res[key] mapping, we need key->token too
        # token_map is token->value; derive key from token "{{key}}"
        for r in runs:
            text = r.text
            if not text:
                continue

            new_text = text
            for token, val in token_map.items():
                if token in new_text:
                    new_text = new_text.replace(token, val)
                    # Extract key name from token "{{key}}"
                    key = token[2:-2]
                    # Store first run where we replaced this key
                    if key not in res:
                        res[key] = r
                    finished.add(key)

            if new_text != text:
                r.text = new_text

    def fill_texts(self, mapping: Dict[str, str]) -> Dict[str, PPTXRun]:
        """
        Replace {{key}} tokens in text frames and table cells.

        Returns key -> PPTXRun where the replacement occurred (best-effort, first run).
        """
        token_map = {f"{{{{{k}}}}}": str(v) for k, v in mapping.items()}

        finished: set[str] = set()
        res: Dict[str, "PPTXRun"] = {}

        for slide in self.presentation.slides:
            for shape in self._iter_shapes(slide.shapes):
                # Text frames
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    tf = shape.text_frame
                    for p in tf.paragraphs:
                        self._replace_tokens_in_runs(p.runs, token_map, res, finished)

                # Tables
                if getattr(shape, "has_table", False) and shape.has_table:
                    tbl = shape.table
                    # python-pptx exposes rows/cols; iterate directly
                    for row in tbl.rows:
                        for cell in row.cells:
                            tf = getattr(cell, "text_frame", None)
                            if tf is None:
                                continue
                            for p in tf.paragraphs:
                                self._replace_tokens_in_runs(p.runs, token_map, res, finished)

        unknown = set(mapping.keys()) - finished
        if unknown:
            _logger.warning(f"Unknown filling keys: {sorted(unknown)}")

        return res

    def save(
            self,
            pptx_save_path: Optional[Union[TypePathLike, IO]] = None,
            pdf_save_path: Optional[TypePathLike] = None,
            image_save_path: Optional[Union[TypePathLike, IO]] = None,
            dpi=200,  # Used for when saving as image
    ):
        if not pptx_save_path and not pdf_save_path and not image_save_path:
            raise ValueError('One path must be specified.')

        if pptx_save_path is not None:
            if os_utils.is_path_like(pptx_save_path):
                pptx_save_path = os_utils.format_path_string(pptx_save_path)
                pptx_save_path.parent.mkdir(parents=True, exist_ok=True)
                self.presentation.save(str(pptx_save_path))
            else:
                # If caller passes an IO, python-pptx can save to it directly.
                self.presentation.save(pptx_save_path)

        if pdf_save_path or image_save_path is not None:
            # In memory temp file
            with tempfile.TemporaryDirectory(dir="/dev/shm", prefix="lo_") as temp_dir_name:
                temp_dir = Path(temp_dir_name)
                if pptx_save_path is not None:
                    temp_pptx_path = pptx_save_path
                else:
                    temp_pptx_path = temp_dir / "temp_ppt.pptx"
                    self.presentation.save(str(temp_pptx_path))

                # Use an isolated LO profile to avoid lock/hang issues
                profile_dir = temp_dir / "profile"
                profile_dir.mkdir(parents=True, exist_ok=True)

                # Inject font substitution rules for this run
                self._write_lo_font_substitution(profile_dir, REG_XCU)

                # Convert PPTX -> PDF
                cmd = [
                    "soffice",  # prefer soffice over libreoffice
                    "--headless",
                    "--nologo",
                    "--nolockcheck",
                    "--nodefault",
                    "--norestore",
                    f"-env:UserInstallation=file://{profile_dir}",
                    "--convert-to", "pdf:impress_pdf_Export",
                    "--outdir", str(temp_dir),
                    str(temp_pptx_path),
                ]
                r = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                if r.returncode != 0:
                    raise RuntimeError(
                        f"LibreOffice failed (code {r.returncode}).\n"
                        f"STDERR:\n{r.stderr}\n"
                        f"STDOUT:\n{r.stdout}"
                    )

                # Output PDF name: LO uses input stem
                temp_pdf_path = temp_dir / f"{temp_pptx_path.stem}.pdf"
                if not temp_pdf_path.exists():
                    raise RuntimeError(f"LibreOffice reported success but PDF not found: {temp_pdf_path}")

                if pdf_save_path is not None:
                    pdf_save_path = os_utils.format_path_string(pdf_save_path)
                    pdf_save_path.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(temp_pdf_path, pdf_save_path)

                if image_save_path is not None:
                    # Requires poppler installed for pdf2image on Linux:
                    # sudo apt-get install -y poppler-utils
                    from pdf2image import convert_from_path
                    pages = convert_from_path(str(temp_pdf_path), dpi=dpi)

                    if os_utils.is_path_like(image_save_path):
                        image_save_path = os_utils.format_path_string(image_save_path)
                        image_save_path.parent.mkdir(parents=True, exist_ok=True)
                    pages[0].save(image_save_path)


    @classmethod
    def _iter_shapes(cls, shapes: PPTXSlideShapes):
        for sh in shapes:
            yield sh
            child = getattr(sh, 'shapes', None)
            if child is not None:
                yield from cls._iter_shapes(child)

    @classmethod
    def _fit_image_into_slide(
            cls,
            slide: PPTXSlide,
            image: npt.NDArray[np.uint8],
            bbox: tuple[PPTXLength, PPTXLength, PPTXLength, PPTXLength],
            fit_mode: TypeFitMode,
    ):
        left, top, w_box, h_box = bbox
        ih, iw = image.shape[:2]

        if iw == 0 or ih == 0:
            raise ValueError(f'Invalid image size: {image.shape}')

        if fit_mode == 'contain':
            scale = min(w_box / iw, h_box / ih)
            width_emu, height_emu = iw * scale, ih * scale
        elif fit_mode == 'cover':
            scale = max(w_box / iw, h_box / ih)
            width_emu, height_emu = iw * scale, ih * scale
        elif fit_mode == 'match_height':
            scale = h_box / ih
            width_emu, height_emu = iw * scale, h_box
        elif fit_mode == 'match_width':
            scale = w_box / iw
            width_emu, height_emu = w_box, ih * scale
        else:
            raise ValueError(f'Invalid fit mode: {fit_mode}')

        x = left + (w_box - int(width_emu)) // 2
        y = top + (h_box - int(height_emu)) // 2
        pic = slide.shapes.add_picture(cls._stream_image(image, extension='.png'), x, y)
        pic.width = int(width_emu)
        pic.height = int(height_emu)
        pic.left, pic.top = x, y
        return pic

    @staticmethod
    def clone_presentation(presentation: PPTXPresentation) -> PPTXPresentation:
        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        return create_presentation(buffer)


    @staticmethod
    def _emu_to_px(emu: int, dpi: int = 96) -> int:
        # EMU_PER_INCH = 914400
        return int(round(emu / 914400 * dpi))

    @staticmethod
    def _stream_image(
            image: npt.NDArray[np.uint8], extension: Optional[str] = None
    ) -> io.BytesIO:
        image_stream = io.BytesIO()
        iio.imwrite(image_stream, image, extension=extension)
        image_stream.seek(0)
        return image_stream

    @staticmethod
    def _write_lo_font_substitution(profile_dir: Path, reg_xcu: str) -> None:
        """
        Create a LibreOffice user profile and write registrymodifications.xcu
        so that font substitution rules apply for this run.
        """
        user_dir = profile_dir / "user"
        user_dir.mkdir(parents=True, exist_ok=True)
        (user_dir / "registrymodifications.xcu").write_text(reg_xcu, encoding="utf-8")


REG_XCU = """
<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry">
  <item oor:path="/org.openoffice.Office.Common/Font/Substitution">
    <prop oor:name="ReplaceTable" oor:op="fuse">
      <value>
        <it>
          <prop oor:name="ReplaceFont"><value>Meiryo UI</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>メイリオ UI</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>Meiryo</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>メイリオ</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>MS Gothic</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>MS Mincho</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Serif CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>Yu Gothic</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>Yu Mincho</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Serif CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>ＭＳ ゴシック</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Sans CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
        <it>
          <prop oor:name="ReplaceFont"><value>ＭＳ 明朝</value></prop>
          <prop oor:name="SubstituteFont"><value>Noto Serif CJK JP</value></prop>
          <prop oor:name="Always"><value>true</value></prop>
        </it>
      </value>
    </prop>
  </item>
</oor:items>
"""