#  Copyright (c) 2025 Daisuke Watanabe,
#  Biomedical Imaging Intelligence Laboratory,
#  Nara Institute of Science and Technology,
#  All rights reserved.

import argparse
import json
from pathlib import Path
from typing import List, Dict, Tuple, Set
import re

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle

from pptx import Presentation
import vtk
from vtk.util import numpy_support as vtknp
import scipy.ndimage as ndi
from PIL import Image

EMU_PER_INCH = 914400


def emu_to_px(emu: int, dpi: int = 96) -> int:
    return int(round(emu / EMU_PER_INCH * dpi))


class PPTFiller:
    def __init__(
            self,
            template_report_path: Path,
            fit_mode: str = 'contain',
            fit_mode_map: dict[str, str] | None = None
    ):
        self.prs = Presentation(str(template_report_path))
        assert fit_mode in {'contain', 'cover', 'match_height', 'match_width'}
        self.fit_mode = fit_mode
        self.fit_mode_map = fit_mode_map or {}

    def _iter_shapes(self, shapes):
        for sh in shapes:
            yield sh
            child = getattr(sh, 'shapes', None)
            if child is not None:
                yield from self._iter_shapes(child)

    def _send_to_back(self, pic_shape):
        el = pic_shape._element
        parent = el.getparent()
        parent.remove(el)
        parent.insert(2, el)

    def _fit_into_shape(self, slide, img_path: Path, box, key: str = ''):
        left, top, w_box, h_box = box
        with Image.open(img_path) as im:
            iw, ih = im.size

        if iw == 0 or ih == 0:
            pic = slide.shapes.add_picture(
                str(img_path), left, top, width=w_box, height=h_box)
            pic.left, pic.top = left, top
            return pic

        def __place(width_emu, height_emu):
            x = left + (w_box - int(width_emu)) // 2
            y = top + (h_box - int(height_emu)) // 2

            pic = slide.shapes.add_picture(str(img_path), x, y)
            pic.width = int(width_emu)
            pic.height = int(height_emu)
            pic.left, pic.top = x, y
            return pic

        mode = self.fit_mode_map.get(key, self.fit_mode)
        if mode == 'contain':
            scale = min(w_box / iw, h_box / ih)
            return __place(iw * scale, ih * scale)
        elif mode == 'cover':
            scale = max(w_box / iw, h_box / ih)
            return __place(iw * scale, ih * scale)
        elif mode == 'match_height':
            scale = h_box / ih
            return __place(iw * scale, h_box)
        else:
            scale = w_box / iw
            return __place(w_box, ih * scale)

    def fill_images(
            self,
            image_map: Dict[str, Path],
            send_to_back_keys: Set[str] | None = None
    ) -> None:
        send_to_back_keys = send_to_back_keys or set()

        for slide in self.prs.slides:
            for shape in self._iter_shapes(slide.shapes):
                name = getattr(shape, 'name', '') or ''
                if 'IMG:' not in name:
                    continue

                key = name.split('IMG:', 1)[1].strip()
                if not key or key not in image_map:
                    continue

                pic = self._fit_into_shape(
                    slide, image_map[key],
                    (shape.left, shape.top, shape.width, shape.height), key)
                if key in send_to_back_keys:
                    self._send_to_back(pic)

    def fill_texts(self, mapping: dict[str, str]) -> None:
        tokens = {f'{{{{{k}}}}}': str(v) for k, v in mapping.items()}

        for slide in self.prs.slides:
            for shape in self._iter_shapes(slide.shapes):
                if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            text = r.text
                            for tk, val in tokens.items():
                                if tk in text:
                                    text = text.replace(tk, val)
                            if text != r.text:
                                r.text = text

                if getattr(shape, 'has_table', False) and shape.has_table:
                    tbl = shape.table
                    for r_i in range(len(tbl.rows)):
                        for c_i in range(len(tbl.columns)):
                            cell = tbl.cell(r_i, c_i)
                            tf = getattr(cell, 'text_frame', None)
                            if tf is None:
                                continue
                            for p in tf.paragraphs:
                                for r in p.runs:
                                    text = r.text
                                    for tk, val in tokens.items():
                                        if tk in text:
                                            text = text.replace(tk, val)
                                    if text != r.text:
                                        r.text = text

    def collect_shape_sizes(self, dpi: int = 96) -> Dict[str, Tuple[int, int]]:
        sizes: Dict[str, Tuple[int, int]] = {}

        for slide in self.prs.slides:
            for shape in self._iter_shapes(slide.shapes):
                name = getattr(shape, 'name', '') or ''
                if 'IMG:' not in name:
                    continue

                key = name.split('IMG:', 1)[1].strip()
                if not key:
                    continue

                w_px = emu_to_px(shape.width, dpi)
                h_px = emu_to_px(shape.height, dpi)
                sizes[key] = (w_px, h_px)

        return sizes

    def save(self, output_report_path: Path) -> Path:
        output_report_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_report_path))
        return output_report_path


class VTKImageHelper:
    @staticmethod
    def read_meta_image(path: Path):
        reader = vtk.vtkMetaImageReader()
        reader.SetFileName(str(path))
        reader.Update()
        return reader

    @staticmethod
    def compute_uniform_scale(image: vtk.vtkImageData, limit: float) -> float:
        x0, x1, y0, y1, z0, z1 = image.GetExtent()
        dims = (x1 - x0 + 1, y1 - y0 + 1, z1 - z0 + 1)
        return min(1.0, limit / dims[0], limit / dims[1], limit / dims[2])

    @staticmethod
    def resample_image(port, scale: float, nearest: bool) -> vtk.vtkImageResample:
        resample = vtk.vtkImageResample()
        resample.SetInputConnection(port)

        if nearest:
            resample.SetInterpolationModeToNearestNeighbor()
        else:
            resample.SetInterpolationModeToLinear()

        resample.SetAxisMagnificationFactor(0, scale)
        resample.SetAxisMagnificationFactor(1, scale)
        resample.SetAxisMagnificationFactor(2, scale)
        resample.Update()
        return resample

    @staticmethod
    def reslice_to_reference(
            moving_port, ref_img: vtk.vtkImageData, nearest: bool) -> vtk.vtkImageReslice:
        reslice = vtk.vtkImageReslice()
        reslice.SetInputConnection(moving_port)
        reslice.SetInformationInput(ref_img)

        if nearest:
            reslice.SetInterpolationModeToNearestNeighbor()
        else:
            reslice.SetInterpolationModeToLinear()

        reslice.SetOutputOrigin(ref_img.GetOrigin())
        reslice.SetOutputSpacing(ref_img.GetSpacing())
        reslice.SetOutputExtent(ref_img.GetExtent())
        reslice.Update()
        return reslice

    @staticmethod
    def clip_voi(port, extent: Tuple[int, int, int, int, int, int]) -> vtk.vtkExtractVOI:
        voi = vtk.vtkExtractVOI()
        voi.SetInputConnection(port)
        voi.SetVOI(*extent)
        voi.Update()
        return voi

    @staticmethod
    def ensure_min_thickness(
            extent: Tuple[int, int, int, int, int, int],
            min_thick: int = 3
    ) -> Tuple[int, int, int, int, int, int]:
        x0, x1, y0, y1, z0, z1 = extent
        if (z1 - z0 + 1) < min_thick:
            mid = (z0 + z1) // 2
            half = max(1, min_thick // 2)
            z0 = mid - half
            z1 = z0 + min_thick - 1

        return x0, x1, y0, y1, z0, z1


class VTKRenderHelper:
    @staticmethod
    def numpy_rgb_to_vtk_image(rgb: np.ndarray) -> vtk.vtkImageData:
        h, w, c = rgb.shape
        img = vtk.vtkImageData()
        img.SetDimensions(w, h, 1)

        vtk_arr = vtknp.numpy_to_vtk(rgb.reshape(-1, c), deep=True, array_type=vtk.VTK_UNSIGNED_CHAR)
        vtk_arr.SetNumberOfComponents(c)
        img.GetPointData().SetScalars(vtk_arr)
        return img

    @staticmethod
    def window_to_png_cropped(ren_win: vtk.vtkRenderWindow, output_path: Path, tight: bool = True) -> Path:
        w2i = vtk.vtkWindowToImageFilter()
        w2i.SetInput(ren_win)
        w2i.ReadFrontBufferOff()
        w2i.SetInputBufferTypeToRGBA()
        w2i.SetScale(1)
        w2i.Update()

        vtk_img = w2i.GetOutput()
        w, h, _ = vtk_img.GetDimensions()
        arr = vtknp.vtk_to_numpy(vtk_img.GetPointData().GetScalars()).reshape(h, w, 4)

        if tight:
            alpha = arr[..., 3]
            mask = alpha < 250
            if mask.any():
                ys = np.where(mask.any(axis=1))[0]
                xs = np.where(mask.any(axis=0))[0]
                y0, y1 = int(ys.min()), int(ys.max()) + 1
                x0, x1 = int(xs.min()), int(xs.max()) + 1
                cropped = arr[y0:y1, x0:x1, :3].astype(np.uint8)
                out_vtk = VTKRenderHelper.numpy_rgb_to_vtk_image(cropped)
            else:
                rgb = arr[..., :3].astype(np.uint8)
                out_vtk = VTKRenderHelper.numpy_rgb_to_vtk_image(rgb)
        else:
            rgb = arr[..., :3].astype(np.uint8)
            out_vtk = VTKRenderHelper.numpy_rgb_to_vtk_image(rgb)

        writer = vtk.vtkPNGWriter()
        writer.SetFileName(str(output_path))
        writer.SetInputData(out_vtk)
        writer.Write()
        return output_path

    @staticmethod
    def make_volume(
            input_port,
            color_tf,
            opacity_tf,
            interpolation: str = 'nearest',
            shade: bool = True,
            ambient: float = 0.6,
            diffuse: float = 0.4,
            specular: float = 0.05,
            specular_power: float = 1.0
    ):
        prop = vtk.vtkVolumeProperty()
        prop.SetColor(color_tf)
        prop.SetScalarOpacity(opacity_tf)
        if interpolation == 'nearest':
            prop.SetInterpolationTypeToNearest()
        else:
            prop.SetInterpolationTypeToLinear()

        if shade:
            prop.ShadeOn()

        prop.SetAmbient(ambient)
        prop.SetDiffuse(diffuse)
        prop.SetSpecular(specular)
        prop.SetSpecularPower(specular_power)

        mapper = vtk.vtkSmartVolumeMapper()
        mapper.SetInputConnection(input_port)
        if hasattr(mapper, 'SetRequestedRenderModeToGPU'):
            mapper.SetRequestedRenderModeToGPU()
        mapper.SetBlendModeToComposite()

        vol = vtk.vtkVolume()
        vol.SetMapper(mapper)
        vol.SetProperty(prop)
        return vol

    @staticmethod
    def new_renderer_window(
            img_w: int,
            img_h: int,
            background: Tuple[float, float, float] = (1.0, 1.0, 1.0)
    ) -> Tuple[vtk.vtkRenderer, vtk.vtkRenderWindow]:
        ren = vtk.vtkRenderer()
        ren.SetBackground(*background)
        win = vtk.vtkRenderWindow()
        win.AddRenderer(ren)
        win.SetSize(img_w, img_h)
        win.OffScreenRenderingOn()

        win.SetAlphaBitPlanes(1)
        win.SetMultiSamples(0)
        ren.SetUseDepthPeeling(True)
        ren.SetMaximumNumberOfPeels(200)
        ren.SetOcclusionRatio(0.1)
        if hasattr(ren, 'SetUseDepthPeelingForVolumes'):
            ren.SetUseDepthPeelingForVolumes(True)

        return ren, win

    @staticmethod
    def position_camera(
            ren: vtk.vtkRenderer,
            view: str,
            center,
            cam_offset: float
    ):
        cam = ren.GetActiveCamera()
        cam.SetFocalPoint(center)
        if view == 'front':
            cam.SetPosition(center[0], center[1] - cam_offset, center[2])
        else:
            cam.SetPosition(center[0], center[1] + cam_offset, center[2])
        cam.SetViewUp(0, 0, 1)


class LabelColorHelper:
    @staticmethod
    def load_label_info(label_info_path: Path) -> Dict[str, List[float]]:
        with open(label_info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return data.get('color', {})

    @staticmethod
    def build_label_rgb_for_names(
            color_tbl: Dict[str, List[float]],
            muscle_names: List[str]
    ) -> Dict[str, Tuple[float, float, float]]:
        out: Dict[str, Tuple[float, float, float]] = {}

        for nm in muscle_names:
            vals = color_tbl.get(nm)
            if isinstance(vals, (list, tuple)) and len(vals) >= 4:
                r, g, b = float(vals[1]), float(vals[2]), float(vals[3])
                out[nm] = (
                    max(0.0, min(1.0, r)),
                    max(0.0, min(1.0, g)),
                    max(0.0, min(1.0, b)),
                )
            else:
                out[nm] = (0.6, 0.6, 0.6)
        return out

    @staticmethod
    def build_color_items(
            color_tbl: Dict[str, List[float]],
            exclude_ids: Set[int],
            allowed_names: Set[str] | None = None
    ) -> List[Tuple[int, float, float, float, float]]:
        items: List[Tuple[int, float, float, float, float]] = []

        for name, vals in color_tbl.items():
            if allowed_names is not None and name not in allowed_names:
                continue
            if not (isinstance(vals, (list, tuple)) and len(vals) >= 5):
                continue
            try:
                lbl_id = int(float(vals[0]))
                r, g, b, a = float(vals[1]), float(vals[2]), float(vals[3]), float(vals[4])
            except Exception:
                continue

            if lbl_id == 0 or lbl_id in exclude_ids:
                continue
            items.append((lbl_id, r, g, b, a))

        items.sort(key=lambda x: x[0])
        return items


class LabelExtentHelper:
    @staticmethod
    def _label_array(img: vtk.vtkImageData) -> Tuple[np.ndarray, Tuple[int, int, int, int, int, int]]:
        arr_vtk = img.GetPointData().GetScalars()
        arr = vtknp.vtk_to_numpy(arr_vtk)

        x0, x1, y0, y1, z0, z1 = img.GetExtent()
        nx, ny, nz = (x1 - x0 + 1, y1 - y0 + 1, z1 - z0 + 1)
        arr = arr.reshape((nz, ny, nx))
        return arr, (x0, x1, y0, y1, z0, z1)

    @staticmethod
    def _longest_true_run(bounds: np.ndarray) -> Tuple[int, int] | None:
        best_len = 0
        best = (None, None)
        cur_len = 0
        start = 0

        for i, val in enumerate(bounds.tolist() + [False]):
            if val and cur_len == 0:
                start = i
            if val:
                cur_len += 1
            else:
                if cur_len > best_len:
                    best_len = cur_len
                    best = (start, i - 1)
                cur_len = 0

        if best[0] is None:
            return None
        return best[0], best[1]

    @staticmethod
    def robust_label_bounds_z(
            img: vtk.vtkImageData,
            label_id: int,
            min_area: int = 200,
            min_voxels: int = 500
    ) -> Tuple[int, int]:
        arr, (x0, x1, y0, y1, z0, z1) = LabelExtentHelper._label_array(img)
        mask = (arr == label_id)
        if mask.sum() < min_voxels:
            return z0, z1

        counts = mask.reshape(mask.shape[0], -1).sum(axis=1)
        good = counts >= min_area
        run = LabelExtentHelper._longest_true_run(good)
        if run is not None:
            k0, k1 = run
            return int(k0 + z0), int(k1 + z0)

        s = 4
        down = mask[::s, ::s, ::s]
        lbl, n = ndi.label(down)
        if n == 0:
            return z0, z1

        sizes = ndi.sum(down, lbl, index=range(1, n + 1))
        k = int(1 + np.argmax(sizes))
        zz = np.where(lbl == k)[0]
        if zz.size == 0:
            return z0, z1

        k0 = int(zz.min()) * s + z0
        k1 = int(zz.max()) * s + z0
        k0 = max(z0, min(z1, k0))
        k1 = max(z0, min(z1, k1))
        if k1 < k0:
            k0, k1 = k1, k0

        return k0, k1

    @staticmethod
    def compute_group_extents(label_img: vtk.vtkImageData) -> Dict[str, Tuple[int, int, int, int, int, int]]:
        x0, x1, y0, y1, z0, z1 = label_img.GetExtent()
        pelvis_top, pelvis_bottom = LabelExtentHelper.robust_label_bounds_z(label_img, 1)
        femur_top, femur_bottom = LabelExtentHelper.robust_label_bounds_z(label_img, 2)
        tibia_top, tibia_bottom = LabelExtentHelper.robust_label_bounds_z(label_img, 73)
        psoas_top, psoas_bottom = LabelExtentHelper.robust_label_bounds_z(label_img, 14)

        pelvis_mid = int((pelvis_top + pelvis_bottom) // 2)
        pelvis_1_3 = int((pelvis_bottom + (pelvis_top - pelvis_bottom) / 3))
        femur_mid = int((femur_top + femur_bottom) // 2)
        femur_1_4 = int((femur_top + (femur_bottom - femur_top) / 4))
        tibia_3_4 = int((tibia_bottom + (tibia_top - tibia_bottom) / 4))

        z5_low = min(pelvis_top, psoas_bottom)
        z5_high = max(pelvis_top, psoas_bottom)
        z5_low = max(z0, min(z1, z5_low))
        z5_high = max(z0, min(z1, z5_high))

        return {
            '1': (x0, x1, y0, y1, pelvis_top, z1),
            '2': (x0, x1, y0, y1, femur_top, pelvis_1_3),
            '3': (x0, x1, y0, y1, femur_top, pelvis_mid),
            '4': (x0, x1, y0, y1, z0, tibia_bottom),
            '5': (x0, x1, y0, y1, z5_low, z5_high),
            '6': (x0, x1, y0, y1, femur_mid, pelvis_bottom),
            '7': (x0, x1, y0, y1, pelvis_top, z1),
            '8': (x0, x1, y0, y1, femur_mid, pelvis_bottom),
            '9': (x0, x1, y0, y1, tibia_3_4, pelvis_mid),
            '10': (x0, x1, y0, y1, z0, femur_1_4),
        }


class HUStatsHelper:
    @staticmethod
    def assign_age_group(sex: str, age: float) -> str:
        if sex == '女性':
            for low_age, high_age in [(50, 59), (60, 69), (70, 79), (80, 95)]:
                if low_age <= age <= high_age:
                    return f'{low_age}-{high_age}'
        else:
            for low_age, high_age in [(50, 79), (80, 95)]:
                if low_age <= age <= high_age:
                    return f'{low_age}-{high_age}'

        return 'out-of-range'

    @staticmethod
    def get_mean_std_tables(summary_path: Path, sex: str) -> tuple[pd.DataFrame, pd.DataFrame]:
        mean_df = pd.read_excel(summary_path, sheet_name=f'{sex}_mean', index_col=0)
        std_df = pd.read_excel(summary_path, sheet_name=f'{sex}_std', index_col=0)
        return mean_df, std_df

    @staticmethod
    def compute_global_hu_xlim(
            mean_df: pd.DataFrame,
            std_df: pd.DataFrame,
            row: pd.Series,
            age_group: str,
            muscle_groups: Dict[str, List[str]],
            excludes: Set[str] | None,
            pad_ratio: float = 0.12
    ) -> Tuple[float, float]:
        excludes = excludes or set()
        vals: List[float] = []

        for gnames in muscle_groups.values():
            for name in gnames:
                if name in excludes:
                    continue
                col = f'HU_{name}'
                if age_group in mean_df.index and col in mean_df.columns:
                    m = pd.to_numeric(mean_df.loc[age_group, col], errors='coerce')
                else:
                    m = np.nan
                if age_group in std_df.index and col in std_df.columns:
                    s = pd.to_numeric(std_df.loc[age_group, col], errors='coerce')
                else:
                    s = np.nan
                t = pd.to_numeric(row.get(col), errors='coerce')
                if np.isfinite(m) and np.isfinite(s) and s >= 0.0:
                    vals.append(float(m - s))
                    vals.append(float(m + s))
                if np.isfinite(t):
                    vals.append(float(t))

        if len(vals) == 0:
            return 0.0, 1.0

        x_min = float(np.min(vals))
        x_max = float(np.max(vals))
        if not np.isfinite(x_min) or not np.isfinite(x_max) or x_min == x_max:
            return 0.0, 1.0

        span = x_max - x_min
        pad = max(1e-6, span * float(pad_ratio))
        return x_min - pad, x_max + pad

    @staticmethod
    def define_observations(
            mean_df: pd.DataFrame,
            std_df: pd.DataFrame,
            row: pd.Series,
            age_group: str,
            muscle_groups: Dict[str, List[str]],
            excludes: Set[str] | None = None,
            messages: Dict[int, str] | None = None
    ) -> str:
        excludes = excludes or set()
        out_of_range = 0
        total = 0

        for gnames in muscle_groups.values():
            for name in gnames:
                if name in excludes:
                    continue

                col = f'HU_{name}'
                mean = pd.to_numeric(mean_df.loc[age_group, col], errors='coerce')
                std = pd.to_numeric(std_df.loc[age_group, col], errors='coerce')
                val = pd.to_numeric(row.get(col), errors='coerce')
                total += 1

                if np.isfinite(mean) and np.isfinite(std) and np.isfinite(val) and val < mean - std:
                    out_of_range += 1

        ratio = out_of_range / max(1, total)
        if ratio >= 0.5:
            return messages[3]
        elif ratio >= 0.2:
            return messages[2]
        else:
            return messages[1]

    @staticmethod
    def star_size_points2_from_box_height(box_h_px: int, dpi: int, scale: float = 0.5) -> float:
        h_pt = box_h_px * 72.0 / dpi
        return (h_pt * scale) ** 2

    @staticmethod
    def draw_mean_std_hu_boxes(
            mean_df: pd.DataFrame,
            std_df: pd.DataFrame,
            row: pd.Series,
            patient_id: str,
            group_name: str,
            muscle_names: list[str],
            excludes: Set[str] | None,
            label_info: Dict[str, Tuple[float, float, float]],
            age_group: str,
            output_dir: Path,
            sex: str,
            xlim: Tuple[float, float] | None = None,
            return_path: bool = False,
            canvas_size_px: Tuple[int, int] | None = None,
            dpi: int = 96,
            s_star: float = 400.0,
            box_height_px: int = 14,
            row_gap_px: int = 8
    ):
        excludes = excludes or set()
        display_names = list(muscle_names)

        if canvas_size_px is None:
            y_pos = np.arange(len(display_names))[::-1]
            fig_h = max(3.0, 1.1 * len(display_names) + 0.8)
            fig, ax = plt.subplots(figsize=(9, fig_h), dpi=150)
            pixel_mode = False
        else:
            w_px, h_px = int(canvas_size_px[0]), int(canvas_size_px[1])
            fig = plt.figure(figsize=(w_px / dpi, h_px / dpi), dpi=dpi)
            ax = fig.add_axes([0, 0, 1, 1])
            pixel_mode = True
        if xlim is None:
            x_min, x_max = 0.0, 1.0
        else:
            x_min, x_max = float(xlim[0]), float(xlim[1])
            if not np.isfinite(x_min) or not np.isfinite(x_max) or x_min == x_max:
                x_min, x_max = 0.0, 1.0

        ax.set_xlim(x_min, x_max)
        ax.set_xticks(np.linspace(x_min, x_max, 6))
        ax.set_xticklabels([])

        if pixel_mode:
            pad_px = 6.0
            n = len(display_names)
            avail = max(1.0, h_px - 2 * pad_px)
            total_needed = n * (box_height_px + row_gap_px) - row_gap_px

            if total_needed > avail:
                scale = avail / max(1.0, total_needed)
                box_h = max(6.0, box_height_px * scale)
                gap_h = row_gap_px * scale
            else:
                box_h = float(box_height_px)
                gap_h = float(row_gap_px)

            ax.set_ylim(0, h_px)
            ax.invert_yaxis()

            n = len(display_names)
            box_h = max(6, int(round(box_h)))
            gap_h = int(round(gap_h))

            avail_int = h_px - 2 * int(pad_px)
            used_int = n * box_h + (n - 1) * gap_h
            top_margin = int(pad_px) + max(0, (avail_int - used_int) // 2)
            y_centers = [top_margin + i * (box_h + gap_h) + box_h // 2 for i in range(n)]

            ax.set_ylim(-1, h_px + 1)
            ax.invert_yaxis()
        else:
            y_pos = np.arange(len(display_names))[::-1]
            box_h = 0.7
            y_centers = None

        for i, name in enumerate(display_names):
            if name in excludes:
                continue
            col = f'HU_{name}'
            if (age_group in mean_df.index and col in mean_df.columns
                    and age_group in std_df.index and col in std_df.columns):
                mean = pd.to_numeric(mean_df.loc[age_group, col], errors='coerce')
                std = pd.to_numeric(std_df.loc[age_group, col], errors='coerce')

                if np.isfinite(mean) and np.isfinite(std) and std >= 0.0:
                    mean = float(mean)
                    std = float(std)
                    x0 = mean - std
                    width = 2 * std

                    if pixel_mode:
                        y_c = y_centers[i]
                        y0 = int(y_c - box_h // 2)
                        rect = Rectangle(
                            (x0, y0), width, int(box_h),
                            facecolor=label_info.get(name),
                            edgecolor='black', linewidth=0.1, antialiased=False)

                        ax.add_patch(rect)
                        half = int(round(box_h * 0.48))
                        ax.plot(
                            [mean, mean],
                            [y_c - half, y_c + half],
                            color='black',
                            linewidth=1.2, antialiased=False)
                    else:
                        rect = Rectangle(
                            (x0, y_pos[i] - 0.35),
                            width,0.7,
                            facecolor=label_info.get(name),
                            edgecolor='black',
                            alpha=0.8, linewidth=0.1)

                        ax.add_patch(rect)
                        ax.plot(
                            [mean, mean],
                            [y_pos[i] - 0.38, y_pos[i] + 0.38],
                            color='black',
                            linewidth=1.2)

            target = pd.to_numeric(row.get(col), errors='coerce')
            if np.isfinite(target):
                target = float(target)
                if pixel_mode:
                    y_c = y_centers[i]
                    ax.scatter(
                        [target], [y_c],
                        marker='*', s=s_star,
                        facecolor='black', edgecolors='white',
                        linewidths=0.7, zorder=5, clip_on=False)
                else:
                    ax.scatter(
                        [target], [y_pos[i]],
                        marker='*', s=s_star,
                        facecolor='black', edgecolors='white',
                        linewidths=0.7, zorder=5)

        for i, name in enumerate(display_names):
            if name in excludes:
                if pixel_mode:
                    y_c = y_centers[i]
                    ax.hlines(
                        int(y_c), x_min, x_max,
                        colors='black',
                        linewidth=1.0, zorder=6, antialiased=False)
                else:
                    ax.hlines(
                        y_pos[i], x_min, x_max,
                        colors='black',
                        linewidth=1.0, zorder=6)

        ax.set_yticks([])
        ax.grid(True)
        for spine in ax.spines.values():
            spine.set_visible(False)

        output_path = output_dir / f'{patient_id}_group_{group_name}_{sex}_{age_group}.png'
        if pixel_mode:
            fig.savefig(output_path, dpi=dpi, bbox_inches=None, pad_inches=0.0)
        else:
            fig.tight_layout(pad=0.1)
            fig.savefig(output_path, dpi=160, bbox_inches='tight')

        plt.close(fig)
        if return_path:
            return output_path


class SegRenderer:
    def __init__(
            self,
            muscle_label_path: Path,
            skin_label_path: Path,
            color_items: List[Tuple[int, float, float, float, float]],
            color_tbl: Dict[str, List[float]],
            exclude_ids: Set[int],
            output_dir: Path,
            voxel_limit: int = 2048,
            image_size: Tuple[int, int] = (1000, 2000),
            camera_offset: float = 2500.0,
            tight_crop: bool = False,
            voi_zoom: float = 1.6,
            keep_aspect_for_ppt: bool = True
    ):
        self.muscle_label_path = muscle_label_path
        self.skin_label_path = skin_label_path
        self.color_items = color_items
        self.color_tbl = color_tbl
        self.exclude_ids = exclude_ids
        self.out_dir = output_dir
        self.voxel_limit = voxel_limit
        self.img_w, self.img_h = image_size
        self.cam_offset = camera_offset
        self.tight_crop = tight_crop
        self.voi_zoom = max(1.0, float(voi_zoom))
        self.keep_aspect_for_ppt = bool(keep_aspect_for_ppt)
        self._label_reader = None
        self._skin_reader = None
        self.label_resampled: vtk.vtkImageResample | None = None
        self.skin_resampled: vtk.vtkImageResample | None = None
        self.skin_volume = None
        self.all_labels_volume = None
        self._renderer: vtk.vtkRenderer | None = None
        self._window: vtk.vtkRenderWindow | None = None

    def build(self):
        self._label_reader = VTKImageHelper.read_meta_image(self.muscle_label_path)
        scale = VTKImageHelper.compute_uniform_scale(self._label_reader.GetOutput(), self.voxel_limit)
        self.label_resampled = VTKImageHelper.resample_image(
            self._label_reader.GetOutputPort(), scale, nearest=True)

        self._skin_reader = VTKImageHelper.read_meta_image(self.skin_label_path)
        skin_resampled0 = VTKImageHelper.resample_image(
            self._skin_reader.GetOutputPort(), scale, nearest=False)
        self.skin_resampled = VTKImageHelper.reslice_to_reference(
            skin_resampled0.GetOutputPort(), self.label_resampled.GetOutput(), nearest=False)

        self.all_labels_volume = self._build_label_volume(self.label_resampled, self.color_items)
        self.skin_volume = self._build_skin_volume(self.skin_resampled)
        self._renderer, self._window = VTKRenderHelper.new_renderer_window(
            self.img_w, self.img_h, background=(1.0, 1.0, 1.0))

    @staticmethod
    def _build_skin_volume(skin_resample: vtk.vtkImageResample):
        smooth = vtk.vtkImageGaussianSmooth()
        smooth.SetInputConnection(skin_resample.GetOutputPort())
        smooth.SetStandardDeviations(0.6, 0.6, 0.6)
        smooth.SetRadiusFactors(1.5, 1.5, 1.5)
        smooth.Update()

        cast = vtk.vtkImageShiftScale()
        cast.SetInputConnection(smooth.GetOutputPort())
        cast.SetShift(0.0)
        cast.SetScale(1.0)
        cast.SetOutputScalarTypeToFloat()
        cast.ClampOverflowOn()
        cast.Update()

        ctf = vtk.vtkColorTransferFunction()
        ctf.AddRGBPoint(0.0, 0.95, 0.80, 0.75)
        ctf.AddRGBPoint(1.0, 0.95, 0.80, 0.75)
        otf = vtk.vtkPiecewiseFunction()
        otf.AddPoint(0.40, 0.00)
        otf.AddPoint(0.48, 0.02)
        otf.AddPoint(0.50, 1.20)
        otf.AddPoint(0.52, 0.02)
        otf.AddPoint(0.60, 0.00)

        vol = VTKRenderHelper.make_volume(
            cast.GetOutputPort(),
            ctf, otf,
            interpolation='linear', shade=True,
            ambient=0.3, diffuse=0.6,
            specular=0.05, specular_power=1.0)
        vol.GetProperty().SetScalarOpacityUnitDistance(3.0)
        return vol

    @staticmethod
    def _build_label_volume(label_resampled_or_voi, color_items: List[Tuple[int, float, float, float, float]]):
        ctf = vtk.vtkColorTransferFunction()
        otf = vtk.vtkPiecewiseFunction()
        otf.AddPoint(-1e6, 0.0)
        otf.AddPoint(1e6, 0.0)

        half = 0.5
        for lbl_id, r, g, b, a in color_items:
            v = float(lbl_id)
            ctf.AddRGBPoint(v, r, g, b)
            otf.AddPoint(v - half, 0.0)
            otf.AddPoint(v, a)
            otf.AddPoint(v + half, 0.0)

        vol = VTKRenderHelper.make_volume(
            label_resampled_or_voi.GetOutputPort(),
            ctf, otf,
            interpolation='nearest', shade=True,
            ambient=0.6, diffuse=0.4,
            specular=0.05, specular_power=1.0)

        return vol

    def _set_window_size(self, w: int, h: int):
        self._window.SetSize(w, h)

    def render_full(self, view: str, output_path: Path) -> Path:
        ren, win = self._renderer, self._window
        ren.RemoveAllViewProps()
        ren.AddVolume(self.all_labels_volume)
        ren.AddVolume(self.skin_volume)

        center = self.skin_resampled.GetOutput().GetCenter()
        VTKRenderHelper.position_camera(
            ren, view,
            center=center, cam_offset=self.cam_offset)

        ren.ResetCameraClippingRange()
        win.Render()
        tight = (self.tight_crop and not self.keep_aspect_for_ppt)
        return VTKRenderHelper.window_to_png_cropped(win, output_path, tight=tight)

    def render_voi(
            self,
            extent: Tuple[int, int, int, int, int, int],
            view: str,
            output_path: Path,
            out_size: Tuple[int, int] | None = None,
            allowed_names: Set[str] | None = None
    ) -> Path:
        if out_size is None:
            out_size = (self.img_w, self.img_h)

        prev_size = self._window.GetSize()
        self._set_window_size(out_size[0], out_size[1])

        extent = VTKImageHelper.ensure_min_thickness(extent, 3)
        label_voi = VTKImageHelper.clip_voi(self.label_resampled.GetOutputPort(), extent)
        skin_voi = VTKImageHelper.clip_voi(self.skin_resampled.GetOutputPort(), extent)
        color_items = LabelColorHelper.build_color_items(
            self.color_tbl, self.exclude_ids, allowed_names=allowed_names)

        label_vol = self._build_label_volume(label_voi, color_items)
        skin_vol = self._build_skin_volume(skin_voi)

        ren, win = self._renderer, self._window
        ren.RemoveAllViewProps()
        ren.AddVolume(label_vol)
        ren.AddVolume(skin_vol)

        center = label_voi.GetOutput().GetCenter()
        effective_offset = self.cam_offset / self.voi_zoom
        VTKRenderHelper.position_camera(
            ren, view,
            center=center, cam_offset=effective_offset)

        ren.ResetCameraClippingRange()
        win.Render()
        tight = (self.tight_crop and not self.keep_aspect_for_ppt)
        out = VTKRenderHelper.window_to_png_cropped(win, output_path, tight=tight)
        self._set_window_size(prev_size[0], prev_size[1])
        return out

def clean_age(s: str) -> int:
    s = str(s)
    digits = ''.join(ch for ch in s if ch.isdigit())
    if digits == '':
        return 0
    return int(digits)

def parse_birthdate(raw_val) -> Tuple[str, str, str]:
    raw_val = str(raw_val).strip()
    parsed_date = pd.to_datetime(raw_val, errors='coerce')
    return f'{parsed_date.year:04d}', f'{parsed_date.month:02d}', f'{parsed_date.day:02d}'

def parse_shooting_date(raw_val) -> Tuple[str, str, str]:
    digits_only = re.sub(r'\D', '', str(raw_val))
    digits_padded = digits_only.zfill(8)[-8:]
    return digits_padded[:4], digits_padded[4:6], digits_padded[6:8]

class ReportGenerator:
    def __init__(self, config: dict):
        self.muscle_label_root = config['muscle_label_root']
        self.skin_label_root = config['skin_label_root']
        self.mean_std_hu_info_path = config['mean_std_hu_info_path']
        self.patient_info_path = config['patient_info_path']
        self.label_info_path = config['label_info_path']
        self.template_report_path = config['template_report_path']
        self.output_report_root = config['output_report_root']
        self.exclude_ids = config['exclude_ids']
        self.dataset_name = config['dataset_name']
        self.muscle_groups = config['muscle_groups']
        self.exclude_by_dataset = config['exclude_by_dataset']
        self.observation_messages = config['observation_messages']

        basic_info_df = pd.read_excel(self.patient_info_path, sheet_name='基本情報')
        hu_df = pd.read_excel(self.patient_info_path, sheet_name='Mean_HU')
        study_info_df = pd.read_excel(self.patient_info_path, sheet_name='Study Info')

        merged_df = pd.merge(basic_info_df, hu_df, on='NAIST_CT_ID', how='inner')
        self.merged_df = pd.merge(merged_df, study_info_df, on='NAIST_CT_ID', how='inner')

        self.color_tbl = LabelColorHelper.load_label_info(self.label_info_path)
        muscle_names_all = sorted({k for v in self.muscle_groups.values() for k in v})
        self.label_rgb = LabelColorHelper.build_label_rgb_for_names(self.color_tbl, muscle_names_all)
        self.color_items_for_renderer = LabelColorHelper.build_color_items(
            self.color_tbl, self.exclude_ids, allowed_names=None
        )

        ppt_probe = PPTFiller(self.template_report_path, fit_mode='match_height', fit_mode_map={})
        shape_sizes_px = ppt_probe.collect_shape_sizes(dpi=96)
        group_keys = [f'{i:02d}' for i in range(1, 11)]
        canvas_heights = [shape_sizes_px[k][1] for k in group_keys if k in shape_sizes_px]
        min_canvas_h = min(canvas_heights) if canvas_heights else 400
        max_rows = max(len(v) for v in self.muscle_groups.values())
        self.row_gap_px = 8
        common_box_height_px = max(10, int((min_canvas_h - max_rows * self.row_gap_px) / max_rows))
        self.common_box_height_px = min(common_box_height_px, 20)
        self.dpi_fig = 96
        self.s_star_common = HUStatsHelper.star_size_points2_from_box_height(
            self.common_box_height_px, dpi=self.dpi_fig, scale=1.5
        )
        self.shape_sizes_px = shape_sizes_px

    def _build_placeholders(self, row: pd.Series, age: int, age_group: str) -> Dict[str, str]:
        name = str(row['氏名'])
        birth_year, brith_month, birth_day = parse_birthdate(str(row['生年月日']))
        height = float(row['身長'])
        weight = float(row['体重'])
        shooting_year, shooting_month, shooting_day = parse_shooting_date(str(row['撮影日']))

        return {
            'NAME': name,
            'SEX': str(row['性別']),
            'BIRTH_YEAR': birth_year,
            'BIRTH_MONTH': brith_month,
            'BIRTH_DAY': birth_day,
            'HEIGHT': height,
            'WEIGHT': weight,
            'AGE': age,
            'SHOOTING_YEAR': shooting_year,
            'SHOOTING_MONTH': shooting_month,
            'SHOOTING_DAY': shooting_day,
        }

    def _process_single_patient(self, row: pd.Series):
        patient_raw = int(row['NAIST_CT_ID'])

        if self.dataset_name == 'uzumasa':
            patient_id = f'UZU{patient_raw:05d}'
            print(patient_id)
        else:
            patient_id = str(patient_raw)

        sex = str(row['性別'])
        age = int(clean_age(str(row['撮影時年齢'])))
        age_group = HUStatsHelper.assign_age_group(sex, age)

        mean_hu_df, std_hu_df = HUStatsHelper.get_mean_std_tables(self.mean_std_hu_info_path, sex)
        excludes = self.exclude_by_dataset.get(self.dataset_name.lower(), set())
        x_min, x_max = HUStatsHelper.compute_global_hu_xlim(
            mean_hu_df, std_hu_df, row, age_group, self.muscle_groups, excludes, pad_ratio=0.12
        )
        observation = HUStatsHelper.define_observations(
            mean_hu_df, std_hu_df, row, age_group, self.muscle_groups, excludes,
            messages=self.observation_messages
        )

        placeholders = self._build_placeholders(row, age, age_group)
        placeholders['OBSERVATIONS'] = observation

        muscle_label_path = self.muscle_label_root / f'{patient_id}_CT1-muscles_label.mhd'
        skin_label_path = self.skin_label_root / f'{patient_id}_CT1' / 'pred_label_original_skin.mha'

        import tempfile
        with tempfile.TemporaryDirectory(prefix=f'group_boxes_{patient_id}_') as tmpdir:
            tmpdir_path = Path(tmpdir)
            image_map: Dict[str, Path] = {}

            for idx, gname in enumerate(sorted(self.muscle_groups.keys(), key=lambda x: int(x)), start=1):
                key = f'{idx:02d}'
                canvas_px = self.shape_sizes_px.get(key)

                output_box_path = HUStatsHelper.draw_mean_std_hu_boxes(
                    mean_hu_df, std_hu_df, row, patient_id, gname,
                    self.muscle_groups[gname], excludes, self.label_rgb, age_group,
                    tmpdir_path, sex, xlim=(x_min, x_max),
                    return_path=True, canvas_size_px=canvas_px,
                    dpi=self.dpi_fig, s_star=self.s_star_common,
                    box_height_px=self.common_box_height_px,
                    row_gap_px=self.row_gap_px
                )
                image_map[key] = output_box_path

            if muscle_label_path.exists() and skin_label_path.exists():
                renderer = SegRenderer(
                    muscle_label_path, skin_label_path,
                    self.color_items_for_renderer, self.color_tbl,
                    self.exclude_ids, tmpdir_path,
                    voxel_limit=2048, image_size=(1500, 2000),
                    camera_offset=2500.0, tight_crop=True,
                    voi_zoom=2.3, keep_aspect_for_ppt=True
                )
                renderer.build()

                all_front = renderer.render_full('front', tmpdir_path / 'ALL_front.png')
                all_back = renderer.render_full('back', tmpdir_path / 'ALL_back.png')
                image_map['11'] = all_front
                image_map['12'] = all_back

                groups_z = LabelExtentHelper.compute_group_extents(
                    renderer.label_resampled.GetOutput())
                front_groups = {'1', '2', '3', '4', '5'}
                insert_idx = 15

                for gkey in sorted(self.muscle_groups.keys(), key=lambda x: int(x)):
                    extent = groups_z.get(gkey)
                    if extent is None:
                        insert_idx += 1
                        continue

                    view = 'front' if gkey in front_groups else 'back'
                    allowed = set(self.muscle_groups[gkey]) - set(excludes)
                    if not allowed:
                        insert_idx += 1
                        continue

                    output_rendering_path = renderer.render_voi(
                        extent, view, tmpdir_path / f'group_{gkey}_{view}.png',
                        out_size=(1500, 2000),
                        allowed_names=allowed
                    )
                    image_map[f'{insert_idx}'] = output_rendering_path
                    insert_idx += 1

                if '19' in image_map:
                    image_map['13'] = image_map['19']
                if '20' in image_map:
                    image_map['14'] = image_map['20']

            fit_mode_map = {f'{i:02d}': 'match_height' for i in range(1, 13)}
            fit_mode_map.update({'15': 'match_width', '21': 'match_width'})

            ppt = PPTFiller(self.template_report_path, fit_mode='match_height', fit_mode_map=fit_mode_map)
            ppt.fill_images(image_map, send_to_back_keys={'11', '12'})
            ppt.fill_texts(placeholders)

            output_report_path = self.output_report_root / patient_id / f'{patient_id}_muscle_report.pptx'
            output_report_path.parent.mkdir(parents=True, exist_ok=True)
            ppt.save(output_report_path)

    def run(self):
        for _, row in self.merged_df.iterrows():
            self._process_single_patient(row)

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        '--muscle_label_root',
        type=Path,
        default=Path(r'\\conger\user\Databases\Uzumasa_study\mhd_files\CT_label_msk'),
    )
    parser.add_argument(
        '--skin_label_root',
        type=Path,
        default=Path(r'\\conger\user\Databases\Uzumasa_study\mhd_files\skin_label'),
    )
    parser.add_argument(
        '--mean_std_hu_path',
        type=Path,
        default=Path(r'Y:\watanabe\data\micbon\mean_std_hu.xlsx'),
    )
    parser.add_argument(
        '--patient_info_path',
        type=Path,
        default=Path(r'Y:\watanabe\data\micbon\NAIST_CT_2025_merged data.xlsx'),
    )
    parser.add_argument(
        '--label_info_path',
        type=Path,
        default=Path(r'Y:\watanabe\data\micbon\naist_totalsegmentator_81.json'),
    )
    parser.add_argument(
        '--template_report_path',
        type=Path,
        default=Path(r'Y:\watanabe\data\micbon\MICBON_AI_report_template_p3.pptx'),
    )
    parser.add_argument(
        '--output_report_root',
        type=Path,
        default=Path(r'Y:\watanabe\projects\seg_report_generation_micbon\workspace\bmd_muscle_reports'),
    )
    args = parser.parse_args()

    config = {
        'muscle_label_root': args.muscle_label_root,
        'skin_label_root': args.skin_label_root,
        'mean_std_hu_info_path': args.mean_std_hu_path,
        'patient_info_path': args.patient_info_path,
        'label_info_path': args.label_info_path,
        'template_report_path': args.template_report_path,
        'output_report_root': args.output_report_root,
        'exclude_ids': {0, 41, 42},
        'dataset_name': 'uzumasa',

        'muscle_groups': {
            '1': ['pectoralis_major', 'pectoralis_minor', 'serratus_anterior', 'intercostal_muscles',
                  'rectus_abdominis', 'internal_oblique', 'external_oblique', 'transversus_abdominis'],
            '2': ['adductor_muscles', 'pectineus_muscle', 'gracilis_muscle', 'sartorius_muscle'],
            '3': ['rectus_femoris_muscle',
                  'vastus_lateralis_muscle_and_vastus_intermedius_muscle',
                  'vastus_medialis_muscle'],
            '4': ['anterior_compartment_muscles', 'lateral_compartment_muscles'],
            '5': ['psoas_major_muscle', 'iliacus_muscle', 'quadratus_lumborum'],
            '6': ['obturator_internus_muscle', 'obturator_externus_muscle', 'piriformis_muscle'],
            '7': ['erector_spinae', 'latissimus_dorsi', 'trapezius', 'supraspinatus',
                  'infraspinatus', 'serratus_anterior', 'subscapularis',
                  'teres_minor_muscle', 'teres_minor'],
            '8': ['gluteus_maximus_muscle', 'gluteus_medius_muscle',
                  'gluteus_minimus_muscle', 'tensor_fasciae_latae_muscle'],
            '9': ['biceps_femoris_muscle', 'semitendinosus_muscle', 'semimembranosus_muscle'],
            '10': ['superficial_posterior_compartment_muscles', 'deep_posterior_compartment_muscles'],
        },

        'exclude_by_dataset': {
            'uzumasa': {
                'intercostal_muscles',  'pectoralis_major', 'pectoralis_minor', 'serratus_anterior',
                'trapezius', 'supraspinatus', 'infraspinatus',
                'subscapularis', 'teres_minor_muscle', 'teres_minor'
            }
        },

        'observation_messages': {
            1: 'あなたの筋肉の質は、同性・同年代と比べて良好あるいは標準的な範囲にあります。'
               'ただし、高齢になると平均的な値でも筋力低下や転倒リスクが高まることがありますので、'
               '今後もバランスの良い食事と（できるだけお医者様や専門家指導の下で）適度な運動を続けることが大切です。',
            2: 'あなたの筋肉の質は、同性・同年代と比べるとやや低めの傾向があります。'
               '高齢期には筋肉の質のわずかな低下でも生活機能に影響が出ることがありますので、'
               '（できるだけお医者様や専門家指導の下で）食事・運動を工夫し、'
               '必要に応じて専門家のアドバイスを受けることをおすすめします。',
            3: 'あなたの筋肉の質は、同性・同年代と比べて低い傾向が見られます。'
               '高齢期では平均的な水準でも転倒や要介護のリスクが高まることが知られています。'
               '早めに医療・リハビリ専門家にご相談されることをおすすめします。',
        },
    }
    generator = ReportGenerator(config)
    generator.run()


if __name__ == '__main__':
    main()
